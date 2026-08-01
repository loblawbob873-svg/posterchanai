import base64
import io
import logging
from typing import Optional

logger = logging.getLogger(__name__)

# Register HEIC/HEIF support for iPhone photos
try:
    from pillow_heif import register_heif_opener
    register_heif_opener()
except ImportError:
    pass  # pillow-heif not installed


def extract_image_text(image_base64: str, max_chars: int = 50000) -> Optional[str]:
    """Extract text from an image using OCR (pytesseract)"""
    try:
        import pytesseract
        from PIL import Image, ExifTags
        import shutil
        import os

        # Try to find tesseract executable - check common paths
        tesseract_cmd = None
        common_paths = [
            '/usr/bin/tesseract',
            '/usr/local/bin/tesseract',
            '/opt/homebrew/bin/tesseract',  # macOS Homebrew
            shutil.which('tesseract'),  # Check PATH
        ]
        
        for path in common_paths:
            if path and os.path.exists(path) and os.access(path, os.X_OK):
                tesseract_cmd = path
                break
        
        if not tesseract_cmd:
            logger.warning("Tesseract OCR not found - OCR will be unavailable. Install tesseract-ocr package.")
            return None
        
        # Set tesseract path explicitly
        pytesseract.pytesseract.tesseract_cmd = tesseract_cmd

        image_bytes = base64.b64decode(image_base64)
        image = Image.open(io.BytesIO(image_bytes))

        # Handle EXIF orientation (phone photos are often rotated)
        try:
            for orientation in ExifTags.TAGS.keys():
                if ExifTags.TAGS[orientation] == 'Orientation':
                    break
            exif = image._getexif()
            if exif is not None:
                orientation_value = exif.get(orientation)
                if orientation_value == 3:
                    image = image.rotate(180, expand=True)
                elif orientation_value == 6:
                    image = image.rotate(270, expand=True)
                elif orientation_value == 8:
                    image = image.rotate(90, expand=True)
        except (AttributeError, KeyError, TypeError):
            pass  # No EXIF data

        # Resize very large images to avoid memory issues (max 4000px on longest side)
        max_dimension = 4000
        if max(image.size) > max_dimension:
            ratio = max_dimension / max(image.size)
            new_size = (int(image.size[0] * ratio), int(image.size[1] * ratio))
            image = image.resize(new_size, Image.Resampling.LANCZOS)

        # Convert to RGB (handles RGBA, P, L, CMYK, etc.)
        if image.mode in ('RGBA', 'LA', 'P'):
            # Handle transparency by compositing on white background
            background = Image.new('RGB', image.size, (255, 255, 255))
            if image.mode == 'P':
                image = image.convert('RGBA')
            background.paste(image, mask=image.split()[-1] if image.mode in ('RGBA', 'LA') else None)
            image = background
        elif image.mode != 'RGB':
            image = image.convert('RGB')

        # Save to BMP format (uncompressed, always supported by leptonica)
        import tempfile
        with tempfile.NamedTemporaryFile(prefix='pcai_doc_bmp_', suffix='.bmp', delete=False) as tmp:
            tmp_path = tmp.name
            image.save(tmp_path, format='BMP')

        try:
            # ONE OCR pass across every installed language (eng + Thai/Chinese/Arabic/… when their packs are
            # present). tesseract handles a multi-script lang string fine — it reads plain English, a short
            # English sign, AND non-Latin text all correctly in a single pass — so there's no threshold to
            # tune and no way for a "which pass wins" heuristic to pick garbage. This is the fix for the
            # "translate returned the source language" bug: eng-only OCR read nothing on a Thai image; now
            # the Thai pack is in the pass. If the combined pass errors (e.g. one corrupt/version-mismatched
            # traineddata), fall back to eng-only so OCR still works. 'osd' is orientation data, not a language.
            try:
                langs = [l for l in pytesseract.get_languages(config="") if l != "osd"]
            except Exception:
                langs = ["eng"]
            _pref = ["eng", "tha", "chi_sim", "chi_tra", "jpn", "kor", "ara", "rus", "hin", "spa", "fra", "deu"]
            ordered = [l for l in _pref if l in langs] + [l for l in langs if l not in _pref]
            lang = "+".join(ordered) or "eng"
            try:
                text = pytesseract.image_to_string(tmp_path, lang=lang)
            except pytesseract.TesseractError as e:
                logger.warning(f"Tesseract OCR error (lang={lang}), falling back to eng: {e}")
                text = pytesseract.image_to_string(tmp_path, lang="eng")
        except pytesseract.TesseractNotFoundError:
            logger.warning("Tesseract OCR executable not found - OCR unavailable")
            return None
        except pytesseract.TesseractError as e:
            logger.warning(f"Tesseract OCR error: {e}")
            return None
        except Exception as e:
            logger.warning(f"OCR processing error: {e}")
            return None
        finally:
            # Clean up temp file
            try:
                os.unlink(tmp_path)
            except Exception:
                pass  # Ignore cleanup errors

        if text and text.strip():
            result = text.strip()
            if len(result) > max_chars:
                result = result[:max_chars] + "\n\n[Text truncated...]"
            logger.info(f"OCR extracted {len(result)} characters from image")
            return result

        logger.debug("OCR found no text in image")
        return None
    except ImportError:
        logger.warning("pytesseract not installed - OCR unavailable")
        return None
    except Exception as e:
        logger.error(f"OCR extraction error: {e}", exc_info=True)
        return None


def extract_pdf_text(pdf_base64: str, max_chars: int = 50000) -> Optional[str]:
    """Extract text from a base64-encoded PDF. Falls back to OCR for image-based pages."""
    try:
        import fitz  # PyMuPDF
        pdf_bytes = base64.b64decode(pdf_base64)
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")

        text_parts = []
        ocr_pages = []  # pages with no machine-readable text
        for page_num, page in enumerate(doc, 1):
            text = page.get_text()
            if text.strip():
                text_parts.append(f"--- Page {page_num} ---\n{text}")
            else:
                ocr_pages.append((page_num, page))

        # OCR fallback for image-based pages (scanned PDFs)
        if ocr_pages and not text_parts:
            for page_num, page in ocr_pages[:10]:  # limit to first 10 pages
                try:
                    pix = page.get_pixmap(dpi=150)
                    img_bytes = pix.tobytes("jpeg")
                    img_b64 = base64.b64encode(img_bytes).decode("utf-8")
                    ocr_result = extract_image_text(img_b64)
                    if ocr_result and ocr_result.strip():
                        text_parts.append(f"--- Page {page_num} (OCR) ---\n{ocr_result}")
                except Exception as ocr_err:
                    logger.warning(f"PDF page {page_num} OCR failed: {ocr_err}")

        doc.close()
        full_text = "\n\n".join(text_parts)

        if len(full_text) > max_chars:
            full_text = full_text[:max_chars] + "\n\n[Document truncated...]"

        return full_text if full_text.strip() else None
    except Exception as e:
        logger.error(f"PDF extraction error: {e}")
        return None


def extract_docx_text(docx_base64: str, max_chars: int = 50000) -> Optional[str]:
    """Extract text from a base64-encoded DOCX file"""
    try:
        from docx import Document
        docx_bytes = base64.b64decode(docx_base64)
        doc = Document(io.BytesIO(docx_bytes))

        text_parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)

        # Also get text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    text_parts.append(row_text)

        full_text = "\n".join(text_parts)

        if len(full_text) > max_chars:
            full_text = full_text[:max_chars] + "\n\n[Document truncated...]"

        return full_text if full_text.strip() else None
    except Exception as e:
        logger.error(f"DOCX extraction error: {e}")
        return None


def extract_xlsx_text(xlsx_base64: str, max_chars: int = 50000) -> Optional[str]:
    """Extract text from a base64-encoded XLSX file"""
    try:
        from openpyxl import load_workbook
        xlsx_bytes = base64.b64decode(xlsx_base64)
        wb = load_workbook(io.BytesIO(xlsx_bytes), data_only=True)

        text_parts = []
        current_length = 0
        max_rows_per_sheet = 1000  # Limit rows per sheet to prevent huge extractions
        max_sheets = 10  # Limit number of sheets to process
        
        sheets_processed = 0
        for sheet_name in wb.sheetnames[:max_sheets]:
            if sheets_processed >= max_sheets:
                break
            sheets_processed += 1
            
            sheet = wb[sheet_name]
            sheet_header = f"--- Sheet: {sheet_name} ---\n"
            text_parts.append(sheet_header)
            current_length += len(sheet_header)
            
            rows_processed = 0
            for row in sheet.iter_rows(values_only=True):
                if rows_processed >= max_rows_per_sheet:
                    text_parts.append(f"[... {max_rows_per_sheet} rows shown, more rows in sheet ...]")
                    break
                
                # Only process rows with actual data
                row_values = [str(cell) if cell is not None else "" for cell in row]
                # Filter out empty cells
                non_empty = [v for v in row_values if v.strip()]
                if non_empty:
                    row_text = " | ".join(non_empty)
                    row_line = row_text + "\n"
                    
                    # Check if adding this row would exceed limit
                    if current_length + len(row_line) > max_chars:
                        text_parts.append(f"\n[Document truncated at {current_length:,} characters...]")
                        current_length = max_chars + 1
                        break
                    
                    text_parts.append(row_line)
                    current_length += len(row_line)
                    rows_processed += 1
                
                # Early exit if we've hit the limit
                if current_length >= max_chars:
                    break
            
            if current_length >= max_chars:
                break

        wb.close()
        full_text = "".join(text_parts)

        if len(full_text) > max_chars:
            full_text = full_text[:max_chars] + "\n\n[Document truncated...]"

        return full_text if full_text.strip() else None
    except Exception as e:
        logger.error(f"XLSX extraction error: {e}")
        return None


def extract_pptx_text(pptx_base64: str, max_chars: int = 50000) -> Optional[str]:
    """Extract text from a base64-encoded PPTX file"""
    try:
        from pptx import Presentation
        pptx_bytes = base64.b64decode(pptx_base64)
        prs = Presentation(io.BytesIO(pptx_bytes))

        text_parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)

            if slide_text:
                text_parts.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_text))

        full_text = "\n\n".join(text_parts)

        if len(full_text) > max_chars:
            full_text = full_text[:max_chars] + "\n\n[Document truncated...]"

        return full_text if full_text.strip() else None
    except Exception as e:
        logger.error(f"PPTX extraction error: {e}")
        return None


def merge_pdfs(pdf_bytes_list: list) -> Optional[bytes]:
    """Merge a list of PDF byte strings into a single PDF. Returns merged PDF bytes or None on error."""
    try:
        import fitz
        merged = fitz.open()
        for pdf_bytes in pdf_bytes_list:
            doc = fitz.open(stream=pdf_bytes, filetype="pdf")
            merged.insert_pdf(doc)
            doc.close()
        out = merged.tobytes(garbage=4, deflate=True)
        merged.close()
        return out
    except Exception as e:
        logger.error(f"PDF merge error: {e}")
        return None


def extract_document_text(document_base64: str, max_chars: int = 50000) -> Optional[str]:
    """Try to extract text from an Office document (auto-detect type)"""
    # Try each format - the wrong format will fail quickly
    for extractor in [extract_docx_text, extract_xlsx_text, extract_pptx_text]:
        result = extractor(document_base64, max_chars)
        if result:
            return result
    return None
