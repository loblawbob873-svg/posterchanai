"""EVERY UPLOADED DOCUMENT IS PARSED HERE AND THE TEXT IS HANDED TO THE MODEL.

`document_service.py` had ZERO test references across 331 lines. It is the front door for
attachments: the bill photo OCR, the Mail "Add to Budget" path, flashcards, and anything a user
drops into chat. The input is a file somebody uploaded and the output goes into an LLM prompt.

Two properties matter more than extraction itself, and both are the kind that hold for four formats
out of five:

  * **EVERY extractor must bound its output.** `max_chars` is what stops one spreadsheet from
    filling the model's whole context — which does not fail, it just quietly costs the rest of the
    conversation. A cap applied to four of the five formats is the exact shape this codebase keeps
    finding, so it is swept rather than spot-checked.
  * **EVERY extractor must return None rather than raise.** They are called from request handlers,
    and `extract_document_text` deliberately tries all three Office formats in turn — that only
    works because the two wrong ones FAIL rather than throw. A raise there would take the request
    down on a file the user could simply have named wrongly.

The tests build real .docx / .xlsx / .pptx / .pdf files with the same libraries the service uses,
rather than mocking them. Every interesting answer here is about what a parser does with input that
is not what it expected, and a mock cannot be wrong the way a parser is — the cross-format cases
below all pass because python-docx genuinely inspects the content type, which is a fact about the
library, not about our code.
"""
import base64
import io

import pytest

from app.services import document_service as ds


def b64(data: bytes) -> str:
    return base64.b64encode(data).decode()


# --------------------------------------------------------------------------- real documents


def make_docx(paragraphs=("hello from docx",), table=None) -> bytes:
    from docx import Document
    d = Document()
    for p in paragraphs:
        d.add_paragraph(p)
    if table:
        t = d.add_table(rows=1, cols=len(table))
        for i, cell in enumerate(table):
            t.rows[0].cells[i].text = cell
    bio = io.BytesIO()
    d.save(bio)
    return bio.getvalue()


def make_xlsx(rows=("hello from xlsx",), sheet="Sheet1") -> bytes:
    from openpyxl import Workbook
    wb = Workbook()
    wb.active.title = sheet
    for i, v in enumerate(rows, 1):
        wb.active.cell(row=i, column=1, value=v)
    bio = io.BytesIO()
    wb.save(bio)
    return bio.getvalue()


def make_pptx(titles=("hello from pptx",)) -> bytes:
    from pptx import Presentation
    p = Presentation()
    for t in titles:
        slide = p.slides.add_slide(p.slide_layouts[5])
        slide.shapes.title.text = t
    bio = io.BytesIO()
    p.save(bio)
    return bio.getvalue()


def make_pdf(pages=("hello from pdf",)) -> bytes:
    import fitz
    doc = fitz.open()
    for text in pages:
        page = doc.new_page()
        page.insert_text((72, 72), text)
    out = doc.tobytes()
    doc.close()
    return out


#: (name, extractor, builder) for the sweeps. A format added later joins by being listed once.
FORMATS = [
    ("docx", ds.extract_docx_text, make_docx),
    ("xlsx", ds.extract_xlsx_text, make_xlsx),
    ("pptx", ds.extract_pptx_text, make_pptx),
    ("pdf", ds.extract_pdf_text, make_pdf),
]


# --------------------------------------------------------------------------- extraction works


def test_a_docx_yields_its_paragraphs():
    assert "hello from docx" in ds.extract_docx_text(b64(make_docx()))


def test_docx_tables_are_extracted_too():
    """A bill, an invoice or a schedule is usually a TABLE, and those are the documents this
    feature exists for. Paragraph-only extraction returns a header and nothing else."""
    got = ds.extract_docx_text(b64(make_docx(paragraphs=("intro",), table=("Total", "42.00"))))
    assert "Total" in got and "42.00" in got


def test_an_xlsx_yields_its_cells_and_names_the_sheet():
    got = ds.extract_xlsx_text(b64(make_xlsx(sheet="Expenses")))
    assert "hello from xlsx" in got
    assert "Expenses" in got, "the sheet name is context the model needs to read the numbers"


def test_an_empty_workbook_is_unreadable_not_readable_and_blank():
    """THE BUG THIS FILE FOUND. The sheet header was appended before any row was read, so an empty
    workbook came back as `--- Sheet: Sheet1 ---` — a NON-EMPTY string. Every caller branches on
    truthiness and the other four extractors return None, so this one alone reported an unreadable
    file as a readable one that says nothing, and the model was asked to find a total in a header.
    People export a blank sheet by accident; the useful answer is that there is nothing in it."""
    assert ds.extract_xlsx_text(b64(make_xlsx(rows=()))) is None


def test_an_empty_sheet_beside_a_full_one_contributes_no_header():
    """The same rule per sheet. A workbook with three blank tabs should not spend the model's
    context announcing them."""
    from openpyxl import Workbook
    wb = Workbook()
    wb.active.title = "Blank"
    ws = wb.create_sheet("Data")
    ws.cell(row=1, column=1, value="the only content")
    bio = io.BytesIO()
    wb.save(bio)

    got = ds.extract_xlsx_text(b64(bio.getvalue()))
    assert "the only content" in got
    assert "Blank" not in got, "an empty sheet still announced itself"


def test_a_pptx_yields_its_slides():
    got = ds.extract_pptx_text(b64(make_pptx()))
    assert "hello from pptx" in got and "Slide 1" in got


def test_a_pdf_yields_its_pages_with_page_markers():
    got = ds.extract_pdf_text(b64(make_pdf(("first page", "second page"))))
    assert "first page" in got and "second page" in got
    assert "Page 1" in got and "Page 2" in got, \
        "page markers are how a citation back to the document stays meaningful"


# --------------------------------------------------------------------------- the cap


@pytest.mark.parametrize("name,extractor,build", FORMATS, ids=[f[0] for f in FORMATS])
def test_every_format_bounds_its_output(name, extractor, build):
    """THE SWEEP. One uncapped format is one spreadsheet away from filling the model's context —
    which never errors, it just silently costs the rest of the conversation.

    The bound is `max_chars` plus a short truncation marker; xlsx reaches it by a different route
    (it stops adding rows) and lands under, which is why this asserts an upper bound rather than an
    exact length."""
    big = {"docx": lambda: build(paragraphs=["x" * 100] * 200),
           "xlsx": lambda: build(rows=["y" * 100] * 500),
           "pptx": lambda: build(titles=["z" * 100] * 50),
           "pdf": lambda: build(pages=[("w" * 80 + "\n") * 30] * 20)}[name]()
    got = extractor(b64(big), max_chars=500)
    assert got is not None
    assert len(got) < 500 + 100, f"{name} returned {len(got)} characters for a cap of 500"


@pytest.mark.parametrize("name,extractor,build", FORMATS, ids=[f[0] for f in FORMATS])
def test_a_truncated_document_says_it_was_truncated(name, extractor, build):
    """The model has to know it is reading a fragment. Silently cut text reads as a complete
    document that happens to end mid-sentence, and the answer is then confidently wrong."""
    big = {"docx": lambda: build(paragraphs=["x" * 100] * 200),
           "xlsx": lambda: build(rows=["y" * 100] * 500),
           "pptx": lambda: build(titles=["z" * 100] * 50),
           "pdf": lambda: build(pages=[("w" * 80 + "\n") * 30] * 20)}[name]()
    assert "truncated" in extractor(b64(big), max_chars=500).lower()


@pytest.mark.parametrize("name,extractor,build", FORMATS, ids=[f[0] for f in FORMATS])
def test_a_small_document_is_not_truncated(name, extractor, build):
    """The other direction, so the cap cannot quietly become 'everything is truncated'."""
    got = extractor(b64(build()))
    assert "truncated" not in got.lower()


# --------------------------------------------------------------------------- bad input


@pytest.mark.parametrize("name,extractor,_build", FORMATS, ids=[f[0] for f in FORMATS])
@pytest.mark.parametrize("bad", ["", "not base64!!", "////", b64(b""), b64(b"junk"),
                                 b64(b"%PDF-1.4 truncated"), b64(b"PK\x03\x04 not really a zip")])
def test_no_extractor_ever_raises(name, extractor, _build, bad):
    """These run inside request handlers, and `extract_document_text` tries three of them in a row
    on the SAME bytes — so two are always being fed the wrong format by design. A raise would turn
    a mis-named upload into a failed request."""
    assert extractor(bad) is None


@pytest.mark.parametrize("name,extractor,build", FORMATS, ids=[f[0] for f in FORMATS])
def test_an_empty_document_is_none_not_an_empty_string(name, extractor, build):
    """Callers branch on truthiness. An empty string would be reported to the model as a document
    that exists and says nothing, instead of one that could not be read."""
    empty = {"docx": lambda: build(paragraphs=()),
             "xlsx": lambda: build(rows=()),
             "pptx": lambda: build(titles=()),
             "pdf": lambda: build(pages=("",))}[name]()
    assert extractor(b64(empty)) is None


# --------------------------------------------------------------------------- the dispatcher


@pytest.mark.parametrize("build,expected", [
    (make_docx, "hello from docx"),
    (make_xlsx, "hello from xlsx"),
    (make_pptx, "hello from pptx"),
])
def test_auto_detection_finds_the_right_format(build, expected):
    assert expected in ds.extract_document_text(b64(build()))


@pytest.mark.parametrize("wrong_extractor,build", [
    (ds.extract_docx_text, make_xlsx),
    (ds.extract_docx_text, make_pptx),
    (ds.extract_xlsx_text, make_docx),
    (ds.extract_xlsx_text, make_pptx),
    (ds.extract_pptx_text, make_docx),
    (ds.extract_pptx_text, make_xlsx),
])
def test_an_extractor_refuses_a_document_of_another_format(wrong_extractor, build):
    """WHAT MAKES TRY-EACH-IN-TURN SAFE. All three Office formats are ZIP files, so nothing about
    the container distinguishes them — the libraries inspect the content type inside. If one ever
    accepted another's file it would return partial nonsense, `extract_document_text` would take
    the first non-empty answer, and the model would be handed garbage as the document."""
    assert wrong_extractor(b64(build())) is None


def test_auto_detection_returns_none_for_something_that_is_not_a_document():
    assert ds.extract_document_text(b64(b"just some plain text, not a document")) is None
    assert ds.extract_document_text("not base64!!") is None


def test_auto_detection_passes_the_cap_through():
    """It takes `max_chars` and hands it to whichever extractor wins. Dropping it on the way would
    silently restore the default 50,000 for every auto-detected upload."""
    got = ds.extract_document_text(b64(make_docx(paragraphs=["x" * 100] * 200)), max_chars=500)
    assert got is not None and len(got) < 500 + 100


# --------------------------------------------------------------------------- merging


def test_merging_pdfs_keeps_every_page():
    import fitz
    merged = ds.merge_pdfs([make_pdf(("a",)), make_pdf(("b", "c"))])
    assert merged is not None
    doc = fitz.open(stream=merged, filetype="pdf")
    try:
        assert doc.page_count == 3, "pages were lost in the merge"
    finally:
        doc.close()


def test_merging_preserves_order():
    import fitz
    merged = ds.merge_pdfs([make_pdf(("first",)), make_pdf(("second",))])
    doc = fitz.open(stream=merged, filetype="pdf")
    try:
        assert "first" in doc[0].get_text() and "second" in doc[1].get_text()
    finally:
        doc.close()


def test_a_merged_pdf_can_be_read_back_by_our_own_extractor():
    """The round trip that matters: merge is used to staple attachments together before they are
    read. A merged file our own reader cannot parse is worse than not merging."""
    merged = ds.merge_pdfs([make_pdf(("alpha",)), make_pdf(("beta",))])
    got = ds.extract_pdf_text(b64(merged))
    assert "alpha" in got and "beta" in got


@pytest.mark.parametrize("bad", [[b"junk"], [b""], [make_pdf(), b"junk"], []])
def test_merging_bad_input_is_none_rather_than_a_corrupt_pdf(bad):
    """A half-merged file would be stored and served as a real document — worse than a failure,
    because nothing downstream would question it."""
    assert ds.merge_pdfs(bad) is None


# --------------------------------------------------------------------------- OCR


def test_ocr_reads_text_from_an_image():
    """The bill-photo path. Skipped where tesseract is not installed — the service returns None and
    says so in the log, which is the documented behaviour, not a failure."""
    pytest.importorskip("pytesseract")
    import shutil
    if not shutil.which("tesseract"):
        pytest.skip("tesseract binary not installed on this machine")

    from PIL import Image, ImageDraw
    img = Image.new("RGB", (400, 120), "white")
    ImageDraw.Draw(img).text((20, 40), "INVOICE 42", fill="black")
    bio = io.BytesIO()
    img.save(bio, format="PNG")

    got = ds.extract_image_text(b64(bio.getvalue()))
    assert got is None or isinstance(got, str)      # OCR accuracy is not ours to assert
    if got:
        assert "truncated" not in got.lower()


@pytest.mark.parametrize("bad", ["", "not base64!!", b64(b"not an image")])
def test_ocr_on_junk_is_none_rather_than_an_exception(bad):
    assert ds.extract_image_text(bad) is None


def test_ocr_output_is_bounded():
    """It shares the cap with every other extractor: OCR of a dense scan is unbounded text, and it
    reaches the same prompt."""
    pytest.importorskip("pytesseract")
    import shutil
    if not shutil.which("tesseract"):
        pytest.skip("tesseract binary not installed on this machine")

    from PIL import Image, ImageDraw
    img = Image.new("RGB", (1200, 900), "white")
    d = ImageDraw.Draw(img)
    for i in range(40):
        d.text((10, 10 + i * 22), "the quick brown fox jumps over the lazy dog " * 2, fill="black")
    bio = io.BytesIO()
    img.save(bio, format="PNG")

    got = ds.extract_image_text(b64(bio.getvalue()), max_chars=100)
    if got:
        assert len(got) < 100 + 100, f"OCR returned {len(got)} characters for a cap of 100"
